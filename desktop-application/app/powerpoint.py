from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import pandas as pd
import os
import graphics
import csv

imgpath = "./desktop-application/app/graphics/"
fldrList = ['clustertables', 'dials', 'gradient', 'participation', 'questiongraphs', 'questiontables', 'wordchart', 'wordgraphs', 'wordtables']
typList = ['overall', 'department', 'position']
prs = Presentation("./desktop-application/app/powerpoint/empty.pptx")
#qfile = pd.read_csv("./desktop-application/app/questionList.csv")

pictures = []
#qtg = graphics.qtg
#companyname = graphics.companyname

class PresentationDetails:
    def __init__(self, title, date):
        self.title = title
        self.date = date

class Slide:
    def __init__(self, title):
        self.title = title

class ImageList:
    def __init__(self, folder, images):
        self.folder = folder
        self.images = images

class ImageDetails:
    def __init__(self, path, name, typ):
        self.path = path
        self.name = name
        self.typ = typ

def init_image_list(path, folder, files):
    images = [ImageDetails(os.path.join(path, file), file, file.split('_')[0]) for file in files if file.endswith(".png")]
    return ImageList(folder, images)

def init_pres_images(base_path, folder_list, picture_list):
    for folder in folder_list:
        folder_path = os.path.join(base_path, folder)
        files = [file for file in os.listdir(folder_path) if file.endswith(".png")]
        picture_list.append(init_image_list(folder_path, folder, files))

def add_slide_with_title(prs, layout_index, title_text):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    return slide

def add_image_to_slide(slide, image_details, left, top, width):
    slide.shapes.add_picture(image_details.path, Inches(left), Inches(top), Inches(width))

def add_table_to_slide(slide, typ, name, dialpic, qtg, qfile):
    df = None
    for tb in qtg:
        if typ in tb.name and name in tb.name:
            df = tb
            break

    if df is not None:
        # Keep all columns
        df = df.copy()

        # Store original column names
        qn_col = df.columns[0]
        yes_col = df.columns[1]

        # Add new columns
        df.insert(0, 'Attribute', '')
        df.insert(2, 'Description', '')

        # Populate the new columns using the CSV data
        for index, row in df.iterrows():
            qNum_value = row[qn_col]
            match_row = qfile[qfile['qNum'] == qNum_value]
            if not match_row.empty:
                df.at[index, 'Attribute'] = match_row['qSubCat'].values[0]
                df.at[index, 'Description'] = match_row['descript'].values[0]

        # Reorder columns
        df = df[['Attribute', qn_col, 'Description', yes_col]]

        # Ensure the yes_col is numeric
        df[yes_col] = pd.to_numeric(df[yes_col], errors='coerce')
        
        # Split the dataframe into two based on the 60% threshold
        df_positive = df[(df[yes_col] > 60) & (~df[qn_col].isin([61, 62])) | 
                 (df[yes_col] <= 60) & (df[qn_col].isin([61, 62]))]
        df_improvement = df[(df[yes_col] <= 60) & (~df[qn_col].isin([61, 62])) | 
                    (df[yes_col] > 60) & (df[qn_col].isin([61, 62]))]


       # First, ensure both dataframes are properly sorted by Attribute and then by yes_col
        df_positive = df_positive.sort_values(
            by=['Attribute', yes_col], 
            ascending=[True, True]  # True for Attribute (A-Z), False for yes_col (highest to lowest)
        )

        df_improvement = df_improvement.sort_values(
            by=['Attribute', yes_col], 
            ascending=[True, False]  # True for Attribute (A-Z), True for yes_col (lowest to highest)
        )

        def create_table_slide(df_slice, slide_title):
            rows_per_slide = 13
            total_rows = df_slice.shape[0]
            slide_count = (total_rows + rows_per_slide - 1) // rows_per_slide

            for slide_num in range(slide_count):
                start_row = slide_num * rows_per_slide
                end_row = min((slide_num + 1) * rows_per_slide, total_rows)
                df_subset = df_slice.iloc[start_row:end_row]

                if slide_num == 0:
                    slide = add_slide_with_title(prs, 3, slide_title)
                else:
                    slide = add_slide_with_title(prs, 3, f"{slide_title} (cont.)")

                add_image_to_slide(slide, dialpic, 11.5, 0.025, 1.5)

                rows, cols = df_subset.shape[0] + 1, df_subset.shape[1]
                left, top, width, height = Inches(1), Inches(1.75), Inches(9), Inches(2)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                # Set column widths
                table.columns[0].width = Inches(2)  # Attribute column
                table.columns[1].width = Inches(0.5)  # QN column
                table.columns[2].width = Inches(8)     # Description column
                table.columns[3].width = Inches(0.75)  # %Yes column

                # Set the headers
                headers = ['Attribute', qn_col, 'Description', '%Yes']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = str(header)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(91, 155, 213)  # #5b9bd5
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

                # Set row heights
                for i in range(rows):
                    table.rows[i].height = Inches(0.35)

                # Fill in the table with data
                for row in range(1, rows):
                    for col in range(cols):
                        cell = table.cell(row, col)
                        value = df_subset.iloc[row - 1, col]
                        if col == 3:  # %Yes column
                            cell.text = f"{int(value)}%"  # Convert to integer and add % sign
                        else:
                            cell.text = str(value)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row % 2 == 0 else RGBColor(173, 216, 230)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)

        # Create positive attributes slide(s)
        if not df_positive.empty:
            create_table_slide(df_positive, f"{name} {'Departmental' if typ == 'DEP' else 'Position'} Positive Attributes")

        # Create improvements slide(s)
        if not df_improvement.empty:
            create_table_slide(df_improvement, f"{name} {'Departmental' if typ == 'DEP' else 'Position'} Improvements")

def init_pres_slides(prs, pictures, companyname, qtg, qfile):
    # Add title slide
    slide = add_slide_with_title(prs, 0, "Liberty University")
    slide.placeholders[1].text = "Cultural Assessment Leadership Team Review"

    for folder in pictures:
        if folder.folder == 'participation':
            slide = add_slide_with_title(prs, 3, f"Participation Analysis")
            horiz = 0.5
            for pic in folder.images:
                tempname = pic.name.replace('.png', '')
                typ, name = tempname.split('_')[0], tempname.split('_')[1].replace('barchart', '')
                add_image_to_slide(slide, pic, horiz, 1.5, 5.625)
                horiz += 6.5

    # Add dial layout slides
    for folder in pictures:
        if folder.folder == 'dials':
            slide = add_slide_with_title(prs, 1, "Executive Summary")
            for pic in folder.images:
                positions = {
                    "OVERALL_RFP.png": (2.9, 2, 2.25),
                    "OVERALL_"+ companyname +".png": (5.25, 1.6, 3),
                    "OVERALL_EPS.png": (8.35, 2, 2.25),
                    "OVERALL_CM.png": (3.25, 4.35, 2.25),
                    "OVERALL_SrLdr.png": (5.62, 4.65, 2.25),
                    "OVERALL_LdrSpv.png": (8, 4.35, 2.25),
                }
                if pic.name in positions:
                    add_image_to_slide(slide, pic, *positions[pic.name])

    # Define the desired order
    desired_order = ['RFP', 'EPS', 'CM', 'SrLdr', 'LdrSpv']

    # Create a dictionary to store folders by their type
    folder_dict = {}

    # Organize folders by their type
    for folder in pictures:
        if folder.folder == 'questiongraphs':
            for pic in folder.images:
                tempname = pic.name.replace('.png', '')
                typ, name = tempname.split('_')[0], tempname.split('_')[1].replace('barchart', '')
                if name not in folder_dict:
                    folder_dict[name] = []
                folder_dict[name].append((folder, pic, typ))

    # Process the folders in the desired order
    for name in desired_order:
        if name in folder_dict:
            for folder, pic, typ in folder_dict[name]:
                # Find the corresponding dial picture
                dialpicture = None
                for dfolder in pictures:
                    if dfolder.folder == "dials":
                        for dialpic in dfolder.images:
                            if name in dialpic.name:
                                dialpicture = dialpic
                                break
                        if dialpicture:
                            break
                
                if dialpicture:
                    add_table_to_slide(slide, typ, name, dialpicture, qtg, qfile)

                    slide = add_slide_with_title(prs, 3, f"{name} {'Departmental' if typ == 'DEP' else 'Position'} Analysis")
                    add_image_to_slide(slide, pic, 2.5, 1.5, 8)
                    add_image_to_slide(slide, dialpicture, 11.5, 0.025, 1.5)

                    slide = add_slide_with_title(prs, 3, f"{name} {'Departmental' if typ == 'DEP' else 'Position'} Breakout")
                    add_image_to_slide(slide, dialpicture, 11.5, 0.025, 1.5)
                
                    for tfolder in pictures:
                        if tfolder.folder == "questiontables":
                            for tpic in tfolder.images:
                                if name in tpic.name and typ in tpic.name and 'concat' in tpic.name:
                                    add_image_to_slide(slide, tpic, 2.5, 1.6, 8)

    # Add word graphs slides before word chart
    for folder in pictures:
        if folder.folder == 'wordgraphs':
            for pic in folder.images:
                if pic.name == 'DEP_WordBarchart.png':
                    slide = add_slide_with_title(prs, 3, "Department WordstoGraph")
                    add_image_to_slide(slide, pic, 2.5, 1.5, 7.5)
                elif pic.name == 'POS_WordBarchart.png':
                    slide = add_slide_with_title(prs, 3, "Position WordstoGraph")
                    add_image_to_slide(slide, pic, 2.5, 1.5, 7.5)

    #add gradient for overall to one slide
    for folder in pictures:
        if folder.folder == 'gradient':
            pic = folder.images[0]
            slide = add_slide_with_title(prs, 3, "Overall Word Association Gradient")
            add_image_to_slide(slide, pic, 2.5, 3.5, 7.5)

    # Add word chart and word chart words slides side by side
    wordchart_pairs = {}
    for folder in pictures:
        if folder.folder == 'wordchart':
            for pic in folder.images:
                if 'WordChart_Words' in pic.name:
                    name = pic.name.replace('OVERALL_', '').replace('WordChart_Words.png', '')
                    if name in wordchart_pairs:
                        wordchart_pairs[name]['words'] = pic
                    else:
                        wordchart_pairs[name] = {'words': pic}
                else:
                    name = pic.name.replace('OVERALL_', '').replace('WordChart.png', '')
                    if name in wordchart_pairs:
                        wordchart_pairs[name]['chart'] = pic
                    else:
                        wordchart_pairs[name] = {'chart': pic}

    # Sort wordchart pairs to ensure "Overall" charts come first
    sorted_wordchart_pairs = sorted(wordchart_pairs.items(), key=lambda x: 'Overall' not in x[1]['chart'].name)  

    for name, pair in sorted_wordchart_pairs:
        if 'chart' in pair and 'words' in pair:
            slide = add_slide_with_title(prs, 3, f"{name} Word Association Comparison")
            add_image_to_slide(slide, pair['chart'], 0.5, 1.5, 5.625)
            add_image_to_slide(slide, pair['words'], 6.625, 1.5, 5.625)

    # Add cluster tables slides
    for folder in pictures:
        if folder.folder == 'clustertables':
            for pic in folder.images:
                tempname = pic.name.replace('.png', '')
                typ, name, ident = tempname.split('_')[0], tempname.split('_')[1], tempname.split('_')[2].replace('clustertable', '')
                if ident == 'p':
                    tident = "positive"
                else:
                    tident = "negative"

                slide = add_slide_with_title(prs, 1, f"{name} {'Department' if typ == 'DEP' else 'Position'} Word Analysis (" + tident + ")")
                add_image_to_slide(slide, pic, 3.25, 1.5, 5)

    prs.save("./desktop-application/app/powerpoint/test.pptx")




def run(qtg, companyname, qfile):
    # Initialize image paths
    init_pres_images(imgpath, fldrList, pictures)
    # Create presentation slides
    init_pres_slides(prs, pictures, companyname, qtg, qfile)
    print("Powerpoint Generation Complete")